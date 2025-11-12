# Установка: pip install python-pptx markdownify

import argparse
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from markdownify import markdownify as md

def get_shape_details(shape, group_top=0, group_left=0):
    """
    Рекурсивно извлекает детали из фигуры, включая текст и абсолютные координаты.
    """
    # Рассчитываем абсолютное положение для фигур внутри групп
    absolute_top = shape.top + group_top
    absolute_left = shape.left + group_left

    # --- Обработка таблиц ---
    if shape.has_table:
        table = shape.table
        html_table = "<table>"
        # Добавляем <thead> для лучшего форматирования Markdown, если первая строка выглядит как заголовок
        # Это простая эвристика, можно усложнить
        is_header = True
        for i, row in enumerate(table.rows):
            html_table += "<thead>" if i == 0 and is_header else "<tbody><tr>"
            for cell in row.cells:
                tag = "th" if i == 0 and is_header else "td"
                html_table += f"<{tag}>{cell.text}</{tag}>"
            html_table += "</thead>" if i == 0 and is_header else "</tr></tbody>"
        html_table += "</table>"
        return {
            "type": "table",
            "text": md(html_table, tables=['table']),
            "top": absolute_top
        }

    # --- Обработка текста в фигурах ---
    elif shape.has_text_frame:
        text = shape.text.strip()
        if text:
            return {
                "type": "text",
                "text": text,
                "top": absolute_top
            }
    
    # --- Рекурсивный обход групповых фигур ---
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_details = []
        for sub_shape in shape.shapes:
            # Передаем координаты текущей группы, чтобы рассчитать абсолютные координаты для вложенных фигур
            sub_details = get_shape_details(sub_shape, shape.top + group_top, shape.left + group_left)
            if sub_details:
                group_details.append(sub_details)
        return group_details

    return None

def process_presentation_with_layout(file_path):
    """
    Извлекает текст с учетом расположения элементов на слайде.
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Ошибка при открытии файла PPTX: {e}")
        sys.exit(1)

    final_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        final_text.append(f"=== СЛАЙД {slide_num} ===")
        
        # 1. Извлекаем основной заголовок
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
        
        # 2. Собираем все остальные элементы с их положением
        all_slide_elements = []
        for shape in slide.shapes:
            # Пропускаем заголовок, чтобы не дублировать
            if shape == slide.shapes.title:
                continue
            
            details = get_shape_details(shape)
            if details:
                # get_shape_details может вернуть список (для групп) или словарь (для одиночных фигур)
                if isinstance(details, list):
                    all_slide_elements.extend(details)
                else:
                    all_slide_elements.append(details)
        
        # 3. Сортируем все элементы по их вертикальному положению (сверху вниз)
        all_slide_elements.sort(key=lambda x: x['top'])

        # 4. Собираем финальный текст для слайда в правильном порядке
        if title_text:
            final_text.append(f"ЗАГОЛОВОК: {title_text}")

        content_parts = []
        tables_md = []
        
        for element in all_slide_elements:
            if element['type'] == 'table':
                tables_md.append(element['text'])
            else: # type == 'text'
                content_parts.append(element['text'])

        if content_parts:
            final_text.append("СОДЕРЖАНИЕ:\n" + "\n".join(content_parts))
            
        if tables_md:
            final_text.append("ТАБЛИЦА:\n" + "\n".join(tables_md))

        final_text.append("\n")

    return "\n".join(final_text)

# --- Основной блок выполнения ---
if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Конвертирует PPTX в текст, восстанавливая логическую структуру слайда на основе расположения блоков."
    )
    parser.add_argument(
        "file_path", 
        type=str, 
        help="Путь к PPTX файлу для конвертации."
    )
    
    args = parser.parse_args()
    file_path = args.file_path

    if not os.path.exists(file_path):
        print(f"Ошибка: Файл не найден по пути '{file_path}'")
        sys.exit(1)

    if not file_path.lower().endswith(".pptx"):
        print("Ошибка: Этот скрипт предназначен только для файлов .pptx")
        sys.exit(1)

    print(f"Обработка файла с учетом структуры слайдов: {file_path}...")
    formatted_rag_text = process_presentation_with_layout(file_path)

    output_filename = os.path.splitext(file_path)[0] + "_for_rag_structured.txt"
    with open(output_filename, "w", encoding="utf-8") as f:
        f.write(formatted_rag_text)

    print(f"Готово! Файл '{output_filename}' успешно создан с восстановленной структурой.")
