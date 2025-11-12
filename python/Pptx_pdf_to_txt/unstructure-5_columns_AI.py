# Установка: pip install python-pptx markdownify

import argparse
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from markdownify import markdownify as md

# Порог в пикселях для определения, находятся ли элементы на одной строке.
# Может потребовать подстройки для презентаций с большим межстрочным интервалом.
VERTICAL_PROXIMITY_THRESHOLD = 25 

def get_shape_details(shape, group_top=0, group_left=0):
    """
    Рекурсивно извлекает детали из фигуры, включая текст и абсолютные координаты.
    """
    absolute_top = shape.top + group_top
    absolute_left = shape.left + group_left

    if shape.has_table:
        table = shape.table
        html_table = "<table>"
        # Упрощенная генерация HTML, markdownify сам разберется со структурой
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                html_table += f"<td>{cell.text}</td>"
            html_table += "</tr>"
        html_table += "</table>"
        return {
            "type": "table",
            "text": md(html_table, tables=['table']),
            "top": absolute_top,
            "left": absolute_left
        }

    elif shape.has_text_frame:
        text = shape.text.strip()
        if text:
            return {
                "type": "text",
                "text": text,
                "top": absolute_top,
                "left": absolute_left
            }
    
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_details = []
        for sub_shape in shape.shapes:
            sub_details = get_shape_details(sub_shape, shape.top + group_top, shape.left + group_left)
            if sub_details:
                if isinstance(sub_details, list):
                    group_details.extend(sub_details)
                else:
                    group_details.append(sub_details)
        return group_details

    return None

def process_presentation_with_layout(file_path):
    """
    Извлекает текст с учетом расположения элементов, поддерживая колоночные макеты.
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Ошибка при открытии файла PPTX: {e}")
        sys.exit(1)

    final_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        final_text.append(f"=== СЛАЙД {slide_num} ===")
        
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
        
        all_slide_elements = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            
            details = get_shape_details(shape)
            if details:
                if isinstance(details, list):
                    all_slide_elements.extend(details)
                else:
                    all_slide_elements.append(details)
        
        if not all_slide_elements and not title_text:
            continue # Пропускаем пустые слайды

        # --- НОВАЯ ЛОГИКА СОРТИРОВКИ ---

        # 1. Сортируем все элементы по вертикали (сверху вниз)
        all_slide_elements.sort(key=lambda x: x['top'])

        # 2. Группируем элементы в "строки" на основе их вертикальной близости
        rows = []
        current_row = [all_slide_elements[0]]
        for i in range(1, len(all_slide_elements)):
            prev_element = all_slide_elements[i-1]
            current_element = all_slide_elements[i]
            
            if abs(current_element['top'] - prev_element['top']) <= VERTICAL_PROXIMITY_THRESHOLD:
                current_row.append(current_element)
            else:
                rows.append(current_row)
                current_row = [current_element]
        rows.append(current_row) # Добавляем последнюю строку

        # 3. Сортируем каждую "строку" по горизонтали (слева направо)
        final_ordered_elements = []
        for row in rows:
            sorted_row = sorted(row, key=lambda x: x['left'])
            final_ordered_elements.extend(sorted_row)

        # --- КОНЕЦ НОВОЙ ЛОГИКИ ---

        # 4. Собираем финальный текст для слайда в правильном порядке
        if title_text:
            final_text.append(f"ЗАГОЛОВОК: {title_text}")

        content_parts = []
        tables_md = []
        
        for element in final_ordered_elements:
            if element['type'] == 'table':
                tables_md.append(element['text'])
            else: # type == 'text'
                content_parts.append(element['text'])

        if content_parts:
            final_text.append("СОДЕРЖАНИЕ:\n" + "\n".join(content_parts))
            
        if tables_md:
            final_text.append("ТАБЛИЦА:\n" + "\n".join(tables_md))

        final_text.append("\n")

    return "\n".join(final_text)

# --- Основной блок выполнения ---
if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Конвертирует PPTX в текст, восстанавливая структуру, включая колоночные макеты."
    )
    parser.add_argument(
        "file_path", 
        type=str, 
        help="Путь к PPTX файлу для конвертации."
    )
    
    args = parser.parse_args()
    file_path = args.file_path

    if not os.path.exists(file_path):
        print(f"Ошибка: Файл не найден по пути '{file_path}'")
        sys.exit(1)

    if not file_path.lower().endswith(".pptx"):
        print("Ошибка: Этот скрипт предназначен только для файлов .pptx")
        sys.exit(1)

    print(f"Обработка файла с учетом колоночной структуры: {file_path}...")
    formatted_rag_text = process_presentation_with_layout(file_path)

    output_filename = os.path.splitext(file_path)[0] + "_for_rag_columns.txt"
    with open(output_filename, "w", encoding="utf-8") as f:
        f.write(formatted_rag_text)

    print(f"Готово! Файл '{output_filename}' успешно создан с учетом структуры колонок.")
