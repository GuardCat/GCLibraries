# Установка: pip install python-pptx markdownify

import argparse
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from markdownify import markdownify as md

def process_presentation_with_pptx(file_path):
    """
    Извлекает весь текст из PPTX, включая текст в фигурах,
    используя библиотеку python-pptx.
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Ошибка при открытии файла PPTX: {e}")
        sys.exit(1)

    final_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        final_text.append(f"=== СЛАЙД {slide_num} ===")
        
        title_text = ""
        content_parts = []
        tables_md = []

        # 1. Пытаемся найти заголовок (стандартный placeholder)
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()

        # 2. Итерируемся по всем остальным фигурам на слайде
        for shape in slide.shapes:
            # Пропускаем уже обработанный заголовок
            if shape == slide.shapes.title:
                continue

            # --- Обработка таблиц ---
            if shape.has_table:
                table = shape.table
                html_table = "<table>"
                for row in table.rows:
                    html_table += "<tr>"
                    for cell in row.cells:
                        html_table += f"<td>{cell.text}</td>"
                    html_table += "</tr>"
                html_table += "</table>"
                tables_md.append(md(html_table, tables=['table']))

            # --- Обработка текста в фигурах (включая текстовые блоки) ---
            elif shape.has_text_frame:
                # Проверяем, есть ли в фигуре текст
                text = shape.text.strip()
                if text:
                    content_parts.append(text)
            
            # --- Обработка групповых фигур ---
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Рекурсивно обходим фигуры внутри группы
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        text = sub_shape.text.strip()
                        if text:
                            content_parts.append(text)

        # 3. Собираем текст для текущего слайда
        if title_text:
            final_text.append(f"ЗАГОЛОВОК: {title_text}")

        if content_parts:
            final_text.append("СОДЕРЖАНИЕ:\n" + "\n".join(content_parts))
            
        if tables_md:
            final_text.append("ТАБЛИЦА:\n" + "\n".join(tables_md))

        final_text.append("\n")

    return "\n".join(final_text)


# --- Основной блок выполнения ---
if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Конвертирует PPTX в текст, надежно извлекая весь текст, включая текст в фигурах."
    )
    parser.add_argument(
        "file_path", 
        type=str, 
        help="Путь к PPTX файлу для конвертации."
    )
    
    args = parser.parse_args()
    file_path = args.file_path

    if not os.path.exists(file_path):
        print(f"Ошибка: Файл не найден по пути '{file_path}'")
        sys.exit(1)

    if not file_path.lower().endswith(".pptx"):
        print("Ошибка: Этот скрипт предназначен только для файлов .pptx")
        sys.exit(1)

    print(f"Обработка файла с помощью python-pptx: {file_path}...")
    formatted_rag_text = process_presentation_with_pptx(file_path)

    output_filename = os.path.splitext(file_path)[0] + "_for_rag_full.txt"
    with open(output_filename, "w", encoding="utf-8") as f:
        f.write(formatted_rag_text)

    print(f"Готово! Файл '{output_filename}' успешно создан с полным извлечением текста.")

